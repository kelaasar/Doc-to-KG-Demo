import json
from html import escape
from docx import Document
from pptx import Presentation
import pdfplumber
from bs4 import BeautifulSoup

def table_to_json(table):
    headers = []
    data = []

    # Extract table headers
    for th in table.find_all('th'):
        headers.append(th.get_text(strip=True))

    # Extract table rows
    for tr in table.find_all('tr'):
        row = {}
        for i, td in enumerate(tr.find_all(['td', 'th'])):
            # Use header if available, otherwise use index
            header = headers[i] if i < len(headers) else str(i)
            row[header] = td.get_text(strip=True)
        if row:
            data.append(row)
    
    return json.dumps(data, indent=4)

def replace_tables_with_json(html_content):
    soup = BeautifulSoup(html_content, 'lxml')

    for table in soup.find_all('table'):
        json_obj = table_to_json(table)
        json_tag = soup.new_tag("pre")
        json_tag.string = json_obj
        table.replace_with(json_tag)

    return str(soup)

def pdf_to_html(pdf_path):
    with pdfplumber.open(pdf_path) as pdf:
        html = '<html><body>'
        
        for page in pdf.pages:
            # Extract text from the entire page (no cropping)
            page_text = page.extract_text(x_tolerance=1, y_tolerance=3, layout=True)
            if page_text and not page.extract_tables():
                html += f'<pre>{escape(page_text)}</pre>'
            
            # Extract tables from the page and convert to JSON directly
            for table in page.extract_tables():
                # Build table HTML for conversion
                table_data = [row for row in table]

                # Convert the table data to JSON
                json_obj = json.dumps(table_data, indent=4)
                html += f'<pre>{json_obj}</pre>'
                
        html += '</body></html>'
    
    return html

def docx_to_html(docx_path):
    doc = Document(docx_path)
    html = '<html><body>'
    
    # Keep track of tables to remove
    tables_to_remove = []

    for element in doc.element.body:
        if element.tag.endswith('p'):
            para = element
            html += f'<p>{escape(para.text)}</p>'
        elif element.tag.endswith('tbl'):
            # Extract table data
            table_html = '<table>'
            table = element
            table_data = []

            for row in table.xpath(".//w:tr"):
                cells = []
                for cell in row.xpath(".//w:tc"):
                    cell_text = ''.join([escape(p.text) for p in cell.xpath(".//w:p") if p.text is not None])
                    cells.append(cell_text)
                table_data.append(cells)

            # Add table to remove list
            tables_to_remove.append(table)

            # Convert table data to JSON and insert it into HTML
            json_obj = json.dumps(table_data, indent=4)
            html += f'<pre>{json_obj}</pre>'

    # Remove tables from the document
    for table in tables_to_remove:
        table.getparent().remove(table)

    html += '</body></html>'
    return html

def pptx_to_html(pptx_path):
    prs = Presentation(pptx_path)
    html = '<html><body>'
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                html += f'<p>{escape(shape.text)}</p>'
    
    html += '</body></html>'
    return html

def file_to_html(file_path, html_path):
    if file_path.endswith('.pdf'):
        html = pdf_to_html(file_path)
    elif file_path.endswith('.docx'):
        html = docx_to_html(file_path)
    elif file_path.endswith('.pptx'):
        html = pptx_to_html(file_path)
    else:
        raise ValueError("Unsupported file type. Please provide a .pdf, .docx, or .pptx file.")
    
    # Write HTML to file
    with open(html_path, 'w', encoding='utf-8') as html_file:
        html_file.write(html)
    return html

# Example usage
file_to_html('data/N_PR_8715_0026.pdf', 'data/converted/PDF-converted.html')
#file_to_html('data/N_PR_8715_0026.docx', 'data/converted/DOCX-converted.html')
#file_to_html('data/Chocolate Cake Recipe.pptx', 'data/converted/PPTX-converted.html')

def html_to_text(html_path, txt_path):
    # Read the HTML file
    with open(html_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    
    # Parse the HTML content with BeautifulSoup
    soup = BeautifulSoup(html_content, 'lxml')
    
    # Extract text from all <pre> and <p> tags
    text_content = ""
    for tag in soup.find_all(['pre', 'p']):
        text_content += tag.get_text() + "\n"
    
    # Write the extracted text to a .txt file
    with open(txt_path, 'w', encoding='utf-8') as txt_file:
        txt_file.write(text_content)

# Example usage
html_to_text('data/converted/PDF-converted.html', 'data/converted/converted.txt')






# Create chunks grouped by similarity

from unstructured.partition.text import partition_text
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np

file_path = "data/converted/converted.txt"
max_chars_per_chunk = 5000  # Set a hard-coded max character limit for each chunk
num_clusters = 20  # Desired number of clusters

paragraphs = partition_text(filename=file_path)
paragraph_texts = [p.text for p in paragraphs]
model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
embeddings = model.encode(paragraph_texts)

def merge_clusters(cluster_a, cluster_b):
    merged_paragraphs = cluster_a['paragraphs'] + cluster_b['paragraphs']
    merged_embeddings = cluster_a['embeddings'] + cluster_b['embeddings']
    return {'paragraphs': merged_paragraphs, 'embeddings': merged_embeddings}

def get_cluster_character_count(cluster):
    return sum(len(p) for p in cluster['paragraphs'])

# Initialize each paragraph as its own cluster
grouped_paragraphs = [{'paragraphs': [paragraph], 'embeddings': [embedding]} for paragraph, embedding in zip(paragraph_texts, embeddings)]
current_num_clusters = len(grouped_paragraphs)

while current_num_clusters > num_clusters:
    best_similarity = -1
    best_pair = None
    
    # Calculate similarity between adjacent clusters
    for i in range(current_num_clusters - 1):
        cluster_a = grouped_paragraphs[i]
        cluster_b = grouped_paragraphs[i + 1]
        
        cluster_a_embedding = np.mean(cluster_a['embeddings'], axis=0)
        cluster_b_embedding = np.mean(cluster_b['embeddings'], axis=0)
        
        similarity = cosine_similarity([cluster_a_embedding], [cluster_b_embedding])[0][0]
        
        combined_char_count = get_cluster_character_count(cluster_a) + get_cluster_character_count(cluster_b)
        
        # Only consider merging if the character count is within the limit
        if similarity > best_similarity and combined_char_count <= max_chars_per_chunk:
            best_similarity = similarity
            best_pair = (i, i + 1)
    
    # If no valid pair to merge was found, break the loop
    if best_pair is None:
        break
    
    # Merge the most similar adjacent clusters that respect the character limit
    cluster_a_index, cluster_b_index = best_pair
    merged_cluster = merge_clusters(grouped_paragraphs[cluster_a_index], grouped_paragraphs[cluster_b_index])
    
    # Replace the first cluster with the merged one and remove the second
    grouped_paragraphs[cluster_a_index] = merged_cluster
    del grouped_paragraphs[cluster_b_index]
    
    # Update the current number of clusters
    current_num_clusters -= 1

# Print clusters
for cluster_idx, cluster in enumerate(grouped_paragraphs):
    print(f"\nCluster {cluster_idx}:")
    for paragraph in cluster['paragraphs']:
        print(paragraph)




# Add descriptive labels

import os
from openai import OpenAI
from dotenv import load_dotenv, find_dotenv
import json

# load the .env file
_ = load_dotenv(find_dotenv())
client = OpenAI(
   api_key=os.environ.get("OPENAI_API_KEY")
)

model = "gpt-4o"
temperature = 0

chunks = ['\n'.join(cluster['paragraphs']) for cluster in grouped_paragraphs]

def get_summary(client, model, messages, temperature):
    completion = client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=temperature,
    )
    return completion.choices[0].message.content

chunksDict = {}

for chunk in chunks:
    prompt = f"""#
    For the given chunk, I want you to return a dictionary of two elements, the key being a short descriptive label of the chunk, and the value being the chunk content. 
    Below is the output schema:
    {{
        "description": "short desc", 
        "content": "long string of document chunk content"
    }},
    """
    
    messages = [
        {"role": "system", "content": "You are an AI that generates descriptive labels for text chunks."},
        {"role": "user", "content": chunk},
        {"role": "user", "content": prompt},
    ]
    
    response = get_summary(client, model, messages, temperature)
    
    response_json = response.strip().replace("```json", "").replace("```", "")
    
    try:
        response_dict = json.loads(response_json)  
    except json.JSONDecodeError as e:
        print("Failed to parse response as JSON:", e)
        continue
    
    chunksDict[response_dict["description"]] = response_dict["content"]


for key, value in chunksDict.items():
   print(f"Description: {key}")
   print(f"Content: {value}\n")
